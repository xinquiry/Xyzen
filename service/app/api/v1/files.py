import hashlib
import logging
from io import BytesIO
from typing import Any
from urllib.parse import quote
from uuid import UUID

from fastapi import APIRouter, Body, Depends, File, Form, HTTPException, UploadFile, status
from fastapi.responses import StreamingResponse
from PIL import Image, ImageDraw, ImageFont
from sqlmodel.ext.asyncio.session import AsyncSession

from app.common.code import ErrCode, ErrCodeError, handle_auth_error
from app.core.storage import (
    FileCategory,
    FileScope,
    StorageServiceProto,
    create_quota_service,
    detect_file_category,
    generate_storage_key,
    get_storage_service,
)
from app.infra.database import get_session
from app.middleware.auth import get_current_user
from app.models.file import FileCreate, FileRead, FileReadWithUrl, FileUpdate
from app.repos.file import FileRepository
from app.repos.knowledge_set import KnowledgeSetRepository

logger = logging.getLogger(__name__)

router = APIRouter(tags=["files"])


def calculate_file_hash(file_data: bytes) -> str:
    """Calculate SHA256 hash of file data"""
    return hashlib.sha256(file_data).hexdigest()


@router.post("/upload", response_model=FileReadWithUrl, status_code=status.HTTP_201_CREATED)
async def upload_file(
    file: UploadFile = File(...),
    scope: str = Form(FileScope.PRIVATE),
    category: str | None = Form(None),
    folder_id: UUID | None = Form(None),
    knowledge_set_id: UUID | None = Form(None),
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> FileReadWithUrl:
    """
    Upload a file to object storage.

    Args:
        file: The file to upload (multipart/form-data)
        scope: File scope (public/private/generated), default: private
        category: File category (images/documents/audio/others), auto-detected if not provided
        folder_id: Optional folder ID to organize the file
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service instance (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        FileReadWithUrl: The created file record with download URL

    Raises:
        HTTPException: 400 if file validation fails, 413 if quota exceeded, 500 if upload fails
    """
    try:
        # Validate file
        if not file.filename:
            raise ErrCode.INVALID_REQUEST.with_messages("Filename is required")

        # Read file data
        file_data = await file.read()
        if not file_data:
            raise ErrCode.EMPTY_MESSAGE.with_messages("File is empty")

        file_size = len(file_data)

        # Validate storage quota BEFORE uploading
        quota_service = create_quota_service(db)
        await quota_service.validate_upload(user_id, file_size)

        # Calculate file hash
        file_hash = calculate_file_hash(file_data)

        # Auto-detect category if not provided
        if not category:
            category = detect_file_category(file.filename)

        # Determine content type
        content_type = file.content_type
        if not content_type or content_type == "application/octet-stream":
            import mimetypes

            content_type, _ = mimetypes.guess_type(file.filename)

            # Explicit fallback for markdown
            if not content_type and file.filename.lower().endswith(".md"):
                content_type = "text/markdown"

            if not content_type:
                content_type = "application/octet-stream"

        # Generate unique storage key (always create new file, no deduplication)
        storage_key = generate_storage_key(
            user_id=user_id,
            filename=file.filename,
            scope=FileScope(scope),
            category=FileCategory(category) if category else None,
        )

        # Upload to object storage
        file_repo = FileRepository(db)
        file_stream = BytesIO(file_data)
        await storage.upload_file(
            file_data=file_stream,
            storage_key=storage_key,
            content_type=content_type,
            metadata={"user_id": user_id},
        )

        # Create database record
        file_create = FileCreate(
            user_id=user_id,
            storage_key=storage_key,
            original_filename=file.filename,
            content_type=content_type,
            file_size=file_size,
            scope=scope,
            category=category or FileCategory.OTHER,
            file_hash=file_hash,
            metainfo={"original_content_type": file.content_type},
            folder_id=folder_id,
        )

        file_record = await file_repo.create_file(file_create)

        # Link to knowledge set if provided
        if knowledge_set_id:
            try:
                ks_repo = KnowledgeSetRepository(db)
                await ks_repo.validate_access(user_id, knowledge_set_id)
                await ks_repo.link_file_to_knowledge_set(file_record.id, knowledge_set_id)
            except ValueError as e:
                logger.warning(f"Failed to link file to knowledge set during upload: {e}")
                # Don't fail the whole upload if linking fails due to access/existence

        await db.commit()
        await db.refresh(file_record)

        # Use API download endpoint (consistent with message attachments)
        download_url = f"/xyzen/api/v1/files/{file_record.id}/download"

        logger.info(f"File uploaded successfully: {storage_key} by user {user_id} ({file_size / (1024 * 1024):.2f}MB)")

        return FileReadWithUrl(
            **file_record.model_dump(),
            download_url=download_url,
        )

    except ErrCodeError as e:
        logger.error(f"File upload failed: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error during file upload: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/", response_model=list[FileRead])
async def list_files(
    scope: str | None = None,
    category: str | None = None,
    include_deleted: bool = False,
    limit: int = 100,
    offset: int = 0,
    folder_id: UUID | None = None,
    filter_by_folder: bool = False,
    user_id: str = Depends(get_current_user),
    db: AsyncSession = Depends(get_session),
) -> list[FileRead]:
    """
    List files for the current user with optional filters.

    Args:
        scope: Optional scope filter (public/private/generated)
        category: Optional category filter (images/documents/audio/others)
        include_deleted: Whether to include soft-deleted files, default: False
        limit: Maximum number of files to return, default: 100
        offset: Number of files to skip, default: 0
        folder_id: Optional folder ID to filter by.
        filter_by_folder: If True, filters files by folder_id (even if None).
        user_id: Authenticated user ID (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        list[FileRead]: List of file records
    """
    try:
        file_repo = FileRepository(db)
        files = await file_repo.get_files_by_user(
            user_id=user_id,
            scope=scope,
            category=category,
            include_deleted=include_deleted,
            limit=min(limit, 1000),  # Cap at 1000
            offset=offset,
            folder_id=folder_id,
            use_folder_filter=filter_by_folder,
        )

        return [FileRead(**file.model_dump()) for file in files]

    except Exception as e:
        logger.error(f"Failed to list files for user {user_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/{file_id}", response_model=FileReadWithUrl)
async def get_file(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> FileReadWithUrl:
    """
    Get file information by ID with download URL.

    Args:
        file_id: File UUID
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service instance (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        FileReadWithUrl: File record with download URL

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id and file_record.scope != FileScope.PUBLIC:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        # Generate download URL
        download_url = await storage.generate_download_url(file_record.storage_key, expires_in=3600)

        return FileReadWithUrl(
            **file_record.model_dump(),
            download_url=download_url,
        )

    except ErrCodeError as e:
        logger.error(f"Failed to get file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error getting file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/{file_id}/download")
async def download_file(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    """
    Download file by ID.

    Args:
        file_id: File UUID
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service instance (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        StreamingResponse: File content as streaming response

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id and file_record.scope != FileScope.PUBLIC:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        # Download from storage
        file_stream = BytesIO()
        await storage.download_file(file_record.storage_key, file_stream)
        file_stream.seek(0)

        # Encode filename for Content-Disposition header (RFC 5987)
        # Support both ASCII and UTF-8 filenames for better browser compatibility
        ascii_filename = file_record.original_filename.encode("ascii", "ignore").decode("ascii")
        utf8_filename = quote(file_record.original_filename.encode("utf-8"))

        return StreamingResponse(
            file_stream,
            media_type=file_record.content_type,
            headers={
                "Content-Disposition": f"attachment; filename=\"{ascii_filename}\"; filename*=UTF-8''{utf8_filename}",
                "Content-Length": str(file_record.file_size),
            },
        )

    except ErrCodeError as e:
        logger.error(f"Failed to download file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error downloading file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/{file_id}/docx-to-pdf")
async def convert_docx_to_pdf(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    """
    Convert a Word document (DOCX) to PDF format and return as preview.

    Args:
        file_id: UUID of the Word file to convert
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service for file access
        db: Database session for file metadata

    Returns:
        StreamingResponse containing the PDF file
    """
    try:
        file_repo = FileRepository(db)
        file = await file_repo.get_file_by_id(file_id)

        if not file:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File not found",
            )

        # Verify user owns this file or it's public
        if str(file.user_id) != str(user_id) and file.scope != FileScope.PUBLIC:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied",
            )

        # Check if file is a Word document
        if not (
            file.content_type
            in [
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/octet-stream",
            ]
            or file.original_filename.lower().endswith((".doc", ".docx"))
        ):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File is not a Word document",
            )

        # Download the Word file from storage
        file_buffer = BytesIO()
        try:
            await storage.download_file(file.storage_key, file_buffer)
        except Exception as e:
            logger.error(f"Failed to download file from storage: {e}")
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data not found in storage",
            )

        file_data = file_buffer.getvalue()
        if not file_data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data is empty",
            )

        # Convert DOCX to PDF
        pdf_buffer = convert_docx_to_pdf_bytes(file_data)

        # Generate safe filename
        base_filename = file.original_filename.rsplit(".", 1)[0]
        pdf_filename = f"{base_filename}.pdf"

        from urllib.parse import quote

        filename_ascii = pdf_filename.encode("utf-8").decode("ascii", errors="ignore") or "document.pdf"
        filename_utf8 = quote(pdf_filename.encode("utf-8"), safe="")

        content_disposition = f"inline; filename=\"{filename_ascii}\"; filename*=UTF-8''{filename_utf8}"

        return StreamingResponse(
            iter([pdf_buffer.getvalue()]),
            media_type="application/pdf",
            headers={"Content-Disposition": content_disposition},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to convert DOCX to PDF for file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to convert file to PDF",
        )


@router.get("/{file_id}/xlsx-to-pdf")
async def convert_xlsx_to_pdf(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    """
    Convert an Excel file (XLSX) to PDF format and return as preview.

    Args:
        file_id: UUID of the Excel file to convert
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service for file access
        db: Database session for file metadata

    Returns:
        StreamingResponse containing the PDF file
    """
    try:
        file_repo = FileRepository(db)
        file = await file_repo.get_file_by_id(file_id)

        if not file:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File not found",
            )

        # Verify user owns this file or it's public
        if str(file.user_id) != str(user_id) and file.scope != FileScope.PUBLIC:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied",
            )

        # Check if file is an Excel file
        if not (
            file.content_type
            in [
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/octet-stream",
            ]
            or file.original_filename.lower().endswith((".xls", ".xlsx"))
        ):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File is not an Excel spreadsheet",
            )

        # Download the Excel file from storage
        file_buffer = BytesIO()
        try:
            await storage.download_file(file.storage_key, file_buffer)
        except Exception as e:
            logger.error(f"Failed to download file from storage: {e}")
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data not found in storage",
            )

        file_data = file_buffer.getvalue()
        if not file_data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data is empty",
            )

        # Convert XLSX to PDF
        pdf_buffer = convert_xlsx_to_pdf_bytes(file_data)

        # Generate safe filename
        base_filename = file.original_filename.rsplit(".", 1)[0]
        pdf_filename = f"{base_filename}.pdf"

        from urllib.parse import quote

        filename_ascii = pdf_filename.encode("utf-8").decode("ascii", errors="ignore") or "document.pdf"
        filename_utf8 = quote(pdf_filename.encode("utf-8"), safe="")

        content_disposition = f"inline; filename=\"{filename_ascii}\"; filename*=UTF-8''{filename_utf8}"

        return StreamingResponse(
            iter([pdf_buffer.getvalue()]),
            media_type="application/pdf",
            headers={"Content-Disposition": content_disposition},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to convert XLSX to PDF for file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to convert file to PDF",
        )


@router.get("/{file_id}/pptx-to-pdf")
async def convert_pptx_to_pdf_get(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    """
    Convert a PowerPoint file (PPTX) to PDF format and return as preview.

    Args:
        file_id: UUID of the PowerPoint file to convert
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service for file access
        db: Database session for file metadata

    Returns:
        StreamingResponse containing the PDF file
    """
    try:
        file_repo = FileRepository(db)
        file = await file_repo.get_file_by_id(file_id)

        if not file:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File not found",
            )

        # Verify user owns this file or it's public
        if str(file.user_id) != str(user_id) and file.scope != FileScope.PUBLIC:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied",
            )

        # Check if file is a PowerPoint file
        if not (
            file.content_type
            in [
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ]
            or file.original_filename.lower().endswith((".ppt", ".pptx"))
        ):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File is not a PowerPoint presentation",
            )

        # Download the PowerPoint file from storage
        file_buffer = BytesIO()
        try:
            await storage.download_file(file.storage_key, file_buffer)
        except Exception as e:
            logger.error(f"Failed to download file from storage: {e}")
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data not found in storage",
            )

        file_data = file_buffer.getvalue()
        if not file_data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="File data is empty",
            )

        # Convert PPTX to PDF
        pdf_buffer = convert_pptx_to_pdf_bytes(file_data)

        # Generate safe filename
        base_filename = file.original_filename.rsplit(".", 1)[0]
        pdf_filename = f"{base_filename}.pdf"

        from urllib.parse import quote

        filename_ascii = pdf_filename.encode("utf-8").decode("ascii", errors="ignore") or "document.pdf"
        filename_utf8 = quote(pdf_filename.encode("utf-8"), safe="")

        content_disposition = f"inline; filename=\"{filename_ascii}\"; filename*=UTF-8''{filename_utf8}"

        return StreamingResponse(
            iter([pdf_buffer.getvalue()]),
            media_type="application/pdf",
            headers={"Content-Disposition": content_disposition},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to convert PPTX to PDF for file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to convert file to PDF",
        )

        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/{file_id}/url", response_model=dict[str, str | int])
async def get_file_url(
    file_id: UUID,
    expires_in: int = 3600,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> dict[str, str | int]:
    """
    Get presigned download URL for a file.

    Args:
        file_id: File UUID
        expires_in: URL expiration time in seconds, default: 3600 (1 hour)
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service instance (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        dict: Dictionary with download_url and expires_in

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id and file_record.scope != FileScope.PUBLIC:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        # Generate presigned URL
        download_url = await storage.generate_download_url(file_record.storage_key, expires_in=expires_in)

        return {
            "download_url": download_url,
            "expires_in": expires_in,
            "storage_key": file_record.storage_key,
        }

    except ErrCodeError as e:
        logger.error(f"Failed to get URL for file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error getting URL for file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.patch("/{file_id}", response_model=FileRead)
async def update_file(
    file_id: UUID,
    file_update: FileUpdate,
    user_id: str = Depends(get_current_user),
    db: AsyncSession = Depends(get_session),
) -> FileRead:
    """
    Update file metadata.

    Args:
        file_id: File UUID
        file_update: File update data
        user_id: Authenticated user ID (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        FileRead: Updated file record

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        # Update file
        updated_file = await file_repo.update_file(file_id, file_update)
        if not updated_file:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        await db.commit()
        await db.refresh(updated_file)

        return FileRead(**updated_file.model_dump())

    except ErrCodeError as e:
        logger.error(f"Failed to update file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error updating file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.delete("/{file_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_file(
    file_id: UUID,
    hard_delete: bool = False,
    user_id: str = Depends(get_current_user),
    storage: StorageServiceProto = Depends(get_storage_service),
    db: AsyncSession = Depends(get_session),
) -> None:
    """
    Delete a file (soft delete by default, hard delete if specified).

    Args:
        file_id: File UUID
        hard_delete: If True, permanently delete file and storage, default: False
        user_id: Authenticated user ID (injected by dependency)
        storage: Storage service instance (injected by dependency)
        db: Database session (injected by dependency)

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        if hard_delete:
            # Delete from object storage
            await storage.delete_file(file_record.storage_key)
            # Delete from database
            await file_repo.hard_delete_file(file_id)
            logger.info(f"File {file_id} hard deleted by user {user_id}")
        else:
            # Soft delete
            await file_repo.soft_delete_file(file_id)
            logger.info(f"File {file_id} soft deleted by user {user_id}")

        await db.commit()

    except ErrCodeError as e:
        logger.error(f"Failed to delete file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error deleting file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.post("/{file_id}/restore", response_model=FileRead)
async def restore_file(
    file_id: UUID,
    user_id: str = Depends(get_current_user),
    db: AsyncSession = Depends(get_session),
) -> FileRead:
    """
    Restore a soft-deleted file.

    Args:
        file_id: File UUID
        user_id: Authenticated user ID (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        FileRead: Restored file record

    Raises:
        HTTPException: 404 if file not found, 403 if access denied
    """
    try:
        file_repo = FileRepository(db)
        file_record = await file_repo.get_file_by_id(file_id)

        if not file_record:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        # Check ownership
        if file_record.user_id != user_id:
            raise ErrCode.FILE_ACCESS_DENIED.with_messages("You don't have access to this file")

        # Restore file
        success = await file_repo.restore_file(file_id)
        if not success:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        await db.commit()

        # Fetch updated record
        restored_file = await file_repo.get_file_by_id(file_id)
        if not restored_file:
            raise ErrCode.FILE_NOT_FOUND.with_messages(f"File {file_id} not found")

        logger.info(f"File {file_id} restored by user {user_id}")

        return FileRead(**restored_file.model_dump())

    except ErrCodeError as e:
        logger.error(f"Failed to restore file {file_id}: {e}")
        raise handle_auth_error(e)
    except Exception as e:
        logger.error(f"Unexpected error restoring file {file_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.get("/stats/summary")
async def get_user_storage_stats(
    user_id: str = Depends(get_current_user),
    db: AsyncSession = Depends(get_session),
) -> dict[str, Any]:
    """
    Get storage statistics and quota information for the current user.

    Args:
        user_id: Authenticated user ID (injected by dependency)
        db: Database session (injected by dependency)

    Returns:
        dict: Dictionary with storage usage, quota limits, and file statistics
    """
    try:
        file_repo = FileRepository(db)
        quota_service = create_quota_service(db)

        total_size = await file_repo.get_total_size_by_user(user_id, include_deleted=False)
        total_files = await file_repo.get_file_count_by_user(user_id, include_deleted=False)
        deleted_files = await file_repo.get_file_count_by_user(user_id, include_deleted=True) - total_files

        # Get quota information
        quota_info = await quota_service.get_quota_info(user_id)

        return {
            "total_files": total_files,
            "total_size": total_size,
            "deleted_files": deleted_files,
            "total_size_mb": round(total_size / (1024 * 1024), 2),
            "total_size_gb": round(total_size / (1024 * 1024 * 1024), 2),
            "quota": quota_info,
        }

    except Exception as e:
        logger.error(f"Failed to get storage stats for user {user_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


@router.delete("/bulk", status_code=status.HTTP_204_NO_CONTENT)
async def bulk_delete_files(
    file_ids: list[UUID] = Body(...),
    user_id: str = Depends(get_current_user),
    db: AsyncSession = Depends(get_session),
) -> None:
    """
    Bulk soft delete multiple files.

    Args:
        file_ids: List of file UUIDs to delete
        user_id: Authenticated user ID (injected by dependency)
        db: Database session (injected by dependency)
    """
    try:
        file_repo = FileRepository(db)
        count = await file_repo.bulk_soft_delete_by_user(user_id, file_ids)
        await db.commit()

        logger.info(f"Bulk deleted {count} files for user {user_id}")

    except Exception as e:
        logger.error(f"Failed to bulk delete files for user {user_id}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e),
        )


def render_pptx_table(table: "pptx.table.Table | None", slide_width_pt: float) -> BytesIO:  # type: ignore  # noqa: F821
    """
    Render a PowerPoint table to PNG image with proper formatting.

    Args:
        table: python-pptx table object
        slide_width_pt: Slide width in points for sizing

    Returns:
        BytesIO object containing the PNG image
    """
    if not table:
        return BytesIO()

    # Table parameters
    rows = table.rows  # type: ignore
    cols = table.columns  # type: ignore
    num_rows = len(rows)
    num_cols = len(cols)

    # Calculate cell dimensions
    cell_width = 120
    cell_height = 40
    border_width = 1

    # Calculate table size
    table_width = cell_width * num_cols + border_width * (num_cols + 1)
    table_height = cell_height * num_rows + border_width * (num_rows + 1)

    # Create image
    table_img = Image.new("RGB", (table_width, table_height), (255, 255, 255))
    draw = ImageDraw.Draw(table_img)

    # Load font
    font_obj = None
    try:
        font_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
        font_obj = ImageFont.truetype(font_path, 14)
    except (OSError, TypeError):
        font_obj = None

    # Draw table
    for row_idx, row in enumerate(rows):
        for col_idx, cell in enumerate(row.cells):
            # Calculate cell position
            x0 = col_idx * cell_width + border_width * (col_idx + 1)
            y0 = row_idx * cell_height + border_width * (row_idx + 1)
            x1 = x0 + cell_width
            y1 = y0 + cell_height

            # Draw cell border
            draw.rectangle([x0, y0, x1, y1], outline=(0, 0, 0), width=border_width)

            # Draw cell background (alternate rows for better visibility)
            if row_idx == 0:  # Header row
                draw.rectangle([x0, y0, x1, y1], fill=(200, 200, 255))
            elif row_idx % 2 == 0:
                draw.rectangle([x0, y0, x1, y1], fill=(240, 240, 240))

            # Extract cell text
            cell_text = cell.text.strip()
            if cell_text:
                # Draw text in cell
                text_x = x0 + 5
                text_y = y0 + 5

                if font_obj:
                    draw.text((text_x, text_y), cell_text, fill=(0, 0, 0), font=font_obj)
                else:
                    draw.text((text_x, text_y), cell_text, fill=(0, 0, 0))

    # Convert to PNG bytes
    png_bytes = BytesIO()
    table_img.save(png_bytes, format="PNG")
    png_bytes.seek(0)

    return png_bytes


def _get_cjk_font_path() -> str:
    """
    Find and return path to a CJK font available on the system.
    Tries multiple common font locations.
    """
    import os

    # Common font paths on Debian/Ubuntu systems
    font_candidates = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSans-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/liberation/LiberationSans-Regular.ttf",  # Fallback
    ]

    for font_path in font_candidates:
        if os.path.exists(font_path):
            logger.info(f"Found CJK font at: {font_path}")
            return font_path

    logger.warning("No CJK font found, will use default font")
    return ""


def convert_pptx_to_pdf_bytes(pptx_data: bytes) -> BytesIO:
    """
    Convert PPTX bytes to PDF bytes using python-pptx and pymupdf with proper text handling.

    Args:
        pptx_data: The raw PPTX file bytes

    Returns:
        BytesIO object containing the PDF data
    """
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        import fitz  # pymupdf
    except ImportError:
        try:
            import pymupdf as fitz
        except ImportError:
            logger.error("pymupdf not installed, cannot convert PPTX to PDF")
            raise ImportError("pymupdf is required for PPTX to PDF conversion")

    try:
        # Load the presentation
        prs = Presentation(io.BytesIO(pptx_data))

        if not prs.slides or len(prs.slides) == 0:
            raise ValueError("Presentation has no slides")

        # Create a new PDF document
        pdf_doc = fitz.open()

        # Slide dimensions: convert from EMUs to points
        # PowerPoint uses EMUs (English Metric Units): 914400 EMUs = 1 inch
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        if slide_width and slide_height:
            slide_width_pt = slide_width / 914400 * 72  # Convert to points
            slide_height_pt = slide_height / 914400 * 72
        else:
            # Default slide size (10" x 7.5")
            slide_width_pt = 720
            slide_height_pt = 540

        logger.info(f"Slide dimensions: {slide_width_pt:.1f} x {slide_height_pt:.1f} points")

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides, 1):
            try:
                logger.info(f"Processing slide {slide_idx} with {len(slide.shapes)} shapes")

                # Add a new page with slide dimensions
                page = pdf_doc.new_page(width=slide_width_pt, height=slide_height_pt)

                # First pass: insert images (background layer)
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                # Extract image data
                                if hasattr(shape, "image") and shape.image:  # type: ignore
                                    image_stream = io.BytesIO(shape.image.blob)  # type: ignore
                                    image_rect = fitz.Rect(
                                        shape.left / 914400 * 72,
                                        shape.top / 914400 * 72,
                                        (shape.left + shape.width) / 914400 * 72,
                                        (shape.top + shape.height) / 914400 * 72,
                                    )

                                    logger.debug(f"Inserting image: {image_rect}")
                                    # Insert image into PDF
                                    page.insert_image(image_rect, stream=image_stream)
                            except Exception as img_err:
                                logger.warning(f"Failed to process image in slide {slide_idx}: {img_err}")
                    except Exception as e:
                        logger.debug(f"Error in image processing loop: {e}")

                # Second pass: insert tables and text (foreground layer)
                text_count = 0
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Skip pictures
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            continue

                        # Handle tables
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            try:
                                if hasattr(shape, "table") and shape.table:  # type: ignore
                                    logger.info(f"[Slide {slide_idx}] Found table with {len(shape.table.rows)} rows")  # type: ignore
                                    table_img = render_pptx_table(shape.table, slide_width_pt)  # type: ignore

                                    # Get table position
                                    table_left_pt = shape.left / 914400 * 72
                                    table_top_pt = shape.top / 914400 * 72
                                    table_width_pt = shape.width / 914400 * 72
                                    table_height_pt = shape.height / 914400 * 72

                                    table_rect = fitz.Rect(
                                        max(0, table_left_pt),
                                        max(0, table_top_pt),
                                        min(slide_width_pt, table_left_pt + table_width_pt),
                                        min(slide_height_pt, table_top_pt + table_height_pt),
                                    )

                                    logger.debug(f"[Slide {slide_idx}] Inserting table at {table_rect}")
                                    page.insert_image(table_rect, stream=table_img, pixmap=None)
                                    logger.info(f"[Slide {slide_idx}] Table inserted successfully")
                            except Exception as table_err:
                                logger.error(f"[Slide {slide_idx}] Failed to render table: {table_err}")
                            continue

                        # Handle text boxes and shapes with text
                        if not hasattr(shape, "text_frame"):
                            continue

                        text_frame = shape.text_frame  # type: ignore
                        if not text_frame:
                            continue

                        full_text = text_frame.text.strip()

                        if not full_text:
                            continue

                        text_count += 1
                        logger.info(f"[Slide {slide_idx}] Found text {text_count}: {full_text[:60]}...")

                        # Calculate position based on shape position
                        shape_left_pt = shape.left / 914400 * 72
                        shape_top_pt = shape.top / 914400 * 72
                        shape_width_pt = max(shape.width / 914400 * 72, 50)  # Minimum width
                        shape_height_pt = max(shape.height / 914400 * 72, 30)  # Minimum height

                        # Determine font size based on paragraph
                        font_size = 10

                        # Check first run for font size
                        for paragraph in text_frame.paragraphs:
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                if first_run.font.size:
                                    font_size = max(float(first_run.font.size.pt), 8)
                                break

                        # Detect title (usually at top and larger)
                        if shape_top_pt < slide_height_pt * 0.25:
                            font_size = max(font_size, 13)

                        # Create text rectangle - add padding to ensure text is visible
                        text_rect = fitz.Rect(
                            max(0, shape_left_pt - 2),
                            max(0, shape_top_pt - 2),
                            min(slide_width_pt, shape_left_pt + shape_width_pt + 2),
                            min(slide_height_pt, shape_top_pt + shape_height_pt + 2),
                        )

                        logger.info(f"[Slide {slide_idx}] Inserting text at {text_rect}, fontsize: {font_size}pt")

                        try:
                            # Use Pillow to render text with proper CJK support
                            # Create image for text rendering
                            text_width = int(shape_width_pt * 2)  # 2x DPI for quality
                            text_height = int(shape_height_pt * 2)

                            # Create transparent image for text
                            text_img = Image.new(
                                "RGBA", (max(text_width, 100), max(text_height, 30)), (255, 255, 255, 0)
                            )
                            text_draw = ImageDraw.Draw(text_img)

                            # Try to load CJK font
                            font_obj = None
                            try:
                                font_size_px = int(font_size * 2)  # Convert to pixels at 2x scale
                                font_path = _get_cjk_font_path()
                                if font_path:
                                    font_obj = ImageFont.truetype(font_path, font_size_px)
                                    logger.info(f"[Slide {slide_idx}] Loaded CJK font at {font_size_px}px")
                            except Exception as font_err:
                                logger.warning(
                                    f"[Slide {slide_idx}] Failed to load CJK font: {font_err}, using default"
                                )
                                # Fall back to default font
                                try:
                                    font_obj = ImageFont.load_default()
                                except (OSError, TypeError):
                                    font_obj = None

                            # Draw text on image
                            if font_obj:
                                text_draw.text((5, 5), full_text, fill=(0, 0, 0, 255), font=font_obj)
                            else:
                                text_draw.text((5, 5), full_text, fill=(0, 0, 0, 255))

                            # Crop to content
                            text_img = text_img.convert("RGB")

                            # Save to bytes
                            text_bytes = BytesIO()
                            text_img.save(text_bytes, format="PNG")
                            text_bytes.seek(0)

                            # Insert image into PDF at text position
                            img_rect = fitz.Rect(
                                max(0, shape_left_pt - 2),
                                max(0, shape_top_pt - 2),
                                min(slide_width_pt, shape_left_pt + shape_width_pt + 2),
                                min(slide_height_pt, shape_top_pt + shape_height_pt + 2),
                            )

                            page.insert_image(img_rect, stream=text_bytes, pixmap=None)
                            logger.info(f"[Slide {slide_idx}] Text inserted as image successfully")

                        except Exception as text_err:
                            logger.error(f"[Slide {slide_idx}] Text insertion error: {text_err}")

                    except Exception as shape_err:
                        logger.warning(f"[Slide {slide_idx}] Failed to process shape {shape_idx}: {shape_err}")

                logger.info(f"[Slide {slide_idx}] Total text elements inserted: {text_count}")

            except Exception as slide_err:
                logger.warning(f"Failed to process slide {slide_idx}: {slide_err}")
                # Continue with next slide instead of failing
                continue

        # Save PDF to bytes
        pdf_buffer = BytesIO()
        pdf_doc.save(pdf_buffer, garbage=0)
        pdf_doc.close()
        pdf_buffer.seek(0)

        logger.info(f"Successfully converted PPTX to PDF ({len(prs.slides)} slides)")
        return pdf_buffer

    except ValueError as ve:
        logger.error(f"Invalid PPTX file: {ve}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid PowerPoint file: {str(ve)}",
        )
    except ImportError as ie:
        logger.error(f"Missing dependencies: {ie}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="PDF conversion library not available",
        )
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to convert file: {str(e)}",
        )


def convert_docx_to_pdf_bytes(docx_data: bytes) -> BytesIO:
    """
    Convert DOCX bytes to PDF bytes using docx2pdf or python-docx + reportlab.

    Args:
        docx_data: The raw DOCX file bytes

    Returns:
        BytesIO object containing the PDF data
    """
    import io

    try:
        import fitz  # pymupdf
    except ImportError:
        try:
            import pymupdf as fitz  # noqa: F401
        except ImportError:
            logger.error("pymupdf not installed, cannot convert DOCX to PDF")
            raise ImportError("pymupdf is required for DOCX to PDF conversion")

    try:
        # Try using libreoffice via command line if available
        import subprocess
        import tempfile
        import os

        # Write DOCX to temporary file
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp_docx:
            tmp_docx.write(docx_data)
            tmp_docx_path = tmp_docx.name

        try:
            # Use LibreOffice to convert DOCX to PDF
            tmp_dir = tempfile.gettempdir()
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    tmp_docx_path,
                ],
                capture_output=True,
                timeout=30,
            )

            # Find the generated PDF
            pdf_path = tmp_docx_path.replace(".docx", ".pdf")
            if os.path.exists(pdf_path):
                with open(pdf_path, "rb") as pdf_file:
                    pdf_data = pdf_file.read()

                # Clean up
                os.remove(tmp_docx_path)
                os.remove(pdf_path)

                logger.info("Successfully converted DOCX to PDF using LibreOffice")
                return BytesIO(pdf_data)
        except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
            logger.warning(f"LibreOffice conversion failed, trying fallback: {e}")
            if os.path.exists(tmp_docx_path):
                os.remove(tmp_docx_path)

        # Fallback: Use python-docx + reportlab with CJK font support for basic conversion
        logger.info("Using python-docx + reportlab with CJK font for DOCX to PDF conversion")
        from docx import Document
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        doc = Document(io.BytesIO(docx_data))

        # Try to register CJK font with reportlab
        try:
            font_path = _get_cjk_font_path()
            if font_path:
                pdfmetrics.registerFont(TTFont("CJKFont", font_path))
                font_name = "CJKFont"
                logger.info(f"Registered CJK font for reportlab: {font_path}")
            else:
                font_name = "Helvetica"
                logger.warning("No CJK font found for reportlab, using Helvetica")
        except Exception as font_err:
            logger.warning(f"Failed to register CJK font with reportlab: {font_err}")
            font_name = "Helvetica"

        # Create PDF
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)
        width, height = letter

        y = height - 0.5 * inch
        line_height = 14

        # Set default font (with CJK support if available)
        c.setFont(font_name, 11)

        # Process paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Handle text wrapping
                text = para.text
                available_width = width - 1 * inch

                # Simple character-based wrapping for CJK text
                # Get approximate character width
                test_width = c.stringWidth("测试", font_name, 11)  # Test with CJK characters
                avg_char_width = test_width / 2 if test_width > 0 else 50

                # Calculate approximate chars per line
                chars_per_line = int(available_width / avg_char_width) if avg_char_width > 0 else 40

                # Split text into chunks
                lines = []
                for i in range(0, len(text), max(chars_per_line, 1)):
                    lines.append(text[i : i + chars_per_line])

                for line in lines:
                    if line.strip():
                        c.drawString(0.5 * inch, y, line)
                        y -= line_height

                # New paragraph spacing
                y -= line_height / 2

            # Check if we need a new page
            if y < 0.5 * inch:
                c.showPage()
                y = height - 0.5 * inch

        # Process tables
        for table in doc.tables:
            y -= line_height

            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        # Truncate long text to fit in cell
                        text = cell.text[:50]
                        c.drawString(0.5 * inch, y, text)
                y -= line_height

            if y < 0.5 * inch:
                c.showPage()
                y = height - 0.5 * inch

        c.save()
        pdf_buffer.seek(0)

        logger.info("Successfully converted DOCX to PDF using python-docx + reportlab")
        return pdf_buffer

    except ImportError as e:
        logger.error(f"Required libraries not installed for DOCX conversion: {e}")
        raise
    except Exception as e:
        logger.error(f"Failed to convert DOCX to PDF: {e}")
        raise


def convert_xlsx_to_pdf_bytes(xlsx_data: bytes) -> BytesIO:
    """
    Convert XLSX bytes to PDF bytes using openpyxl + reportlab or LibreOffice.
    Intelligently determines page orientation based on data dimensions.

    Args:
        xlsx_data: The raw XLSX file bytes

    Returns:
        BytesIO object containing the PDF data
    """
    import io

    # First, analyze the data to determine orientation
    from openpyxl import load_workbook as openpyxl_load

    try:
        wb_temp = openpyxl_load(io.BytesIO(xlsx_data))
        ws_temp = wb_temp.active

        if ws_temp is None:
            use_landscape = True
        else:
            # Count non-empty cells
            max_row = 0
            max_col = 0
            for row in ws_temp.iter_rows(values_only=True):
                for col_idx, cell in enumerate(row):
                    if cell is not None:
                        max_col = max(max_col, col_idx + 1)
                max_row += 1

            # Determine orientation: if columns > rows, use landscape
            use_landscape = max_col > max_row
            logger.info(f"Excel data: {max_row} rows, {max_col} cols → {'landscape' if use_landscape else 'portrait'}")

    except Exception as e:
        logger.warning(f"Could not determine orientation, defaulting to landscape: {e}")
        use_landscape = True

    try:
        import fitz  # pymupdf
    except ImportError:
        try:
            import pymupdf as fitz  # noqa: F401
        except ImportError:
            logger.error("pymupdf not installed, cannot convert XLSX to PDF")
            raise ImportError("pymupdf is required for XLSX to PDF conversion")

    try:
        # Try using libreoffice via command line if available
        import subprocess
        import tempfile
        import os

        # Write XLSX to temporary file
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp_xlsx:
            tmp_xlsx.write(xlsx_data)
            tmp_xlsx_path = tmp_xlsx.name

        try:
            # Use LibreOffice to convert XLSX to PDF
            tmp_dir = tempfile.gettempdir()
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    tmp_xlsx_path,
                ],
                capture_output=True,
                timeout=30,
            )

            # Find the generated PDF
            pdf_path = tmp_xlsx_path.replace(".xlsx", ".pdf")
            if os.path.exists(pdf_path):
                with open(pdf_path, "rb") as pdf_file:
                    pdf_data = pdf_file.read()

                # Clean up
                os.remove(tmp_xlsx_path)
                os.remove(pdf_path)

                logger.info("Successfully converted XLSX to PDF using LibreOffice")
                return BytesIO(pdf_data)
        except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
            logger.warning(f"LibreOffice conversion failed, trying fallback: {e}")
            if os.path.exists(tmp_xlsx_path):
                os.remove(tmp_xlsx_path)

        # Fallback: Use openpyxl + reportlab for basic conversion
        logger.info("Using openpyxl + reportlab for XLSX to PDF conversion")
        from openpyxl import load_workbook
        from reportlab.lib.pagesizes import letter, landscape, portrait
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        from reportlab.lib import colors
        from reportlab.platypus import Table, TableStyle  # noqa: F401

        wb = load_workbook(io.BytesIO(xlsx_data))

        # Choose page size based on orientation
        if use_landscape:
            page_size = landscape(letter)
        else:
            page_size = portrait(letter)

        page_width, page_height = page_size

        # Create PDF with smart orientation
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=page_size)

        margin = 0.3 * inch
        available_width = page_width - 2 * margin
        available_height = page_height - 1 * inch

        # Font configuration
        header_font_size = 11
        cell_font_size = 9

        # Process each worksheet
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            if sheet_idx > 0:
                c.showPage()

            ws = wb[sheet_name]

            # Sheet title
            c.setFont("Helvetica-Bold", 14)
            c.drawString(margin, page_height - 0.4 * inch, f"Sheet: {sheet_name}")

            # Get data from worksheet
            data = []
            col_widths = {}

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                row_data = []
                for col_idx, cell_value in enumerate(row):
                    # Convert value to string
                    cell_text = str(cell_value) if cell_value is not None else ""
                    row_data.append(cell_text)

                    # Track column widths (estimate based on content length)
                    text_len = len(cell_text)
                    col_widths[col_idx] = max(col_widths.get(col_idx, 0), text_len)

                if any(row_data):  # Only add non-empty rows
                    data.append(row_data)

            if not data:
                continue

            # Calculate actual column widths
            num_cols = max(len(row) for row in data)

            # Auto-calculate column widths based on content
            adjusted_widths = []
            total_content_width = 0

            for col_idx in range(num_cols):
                content_width = col_widths.get(col_idx, 5)
                # Convert character count to points (roughly 6 points per character)
                char_width = max(content_width * 6, 30)
                adjusted_widths.append(char_width)
                total_content_width += char_width

            # Scale widths to fit available width
            scale_factor = available_width / max(total_content_width, available_width)
            scaled_widths = [w * scale_factor for w in adjusted_widths]

            # Create table with proper styling
            try:
                table = Table(data, colWidths=scaled_widths)
                table.setStyle(
                    TableStyle(
                        [
                            # Header styling
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D3D3D3")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#000000")),
                            # Cell alignment and padding
                            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                            ("LEFTPADDING", (0, 0), (-1, -1), 4),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                            ("TOPPADDING", (0, 0), (-1, -1), 4),
                            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                            # Font styling
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                            ("FONTSIZE", (0, 0), (-1, 0), header_font_size),
                            ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                            ("FONTSIZE", (0, 1), (-1, -1), cell_font_size),
                            # Grid and borders - IMPORTANT: thick lines for clarity
                            ("GRID", (0, 0), (-1, -1), 1.0, colors.HexColor("#666666")),  # 1pt grid lines
                            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8F8F8")]),
                            # Top border (thick)
                            ("LINEABOVE", (0, 0), (-1, 0), 1.5, colors.HexColor("#000000")),
                            # Bottom border (thick)
                            ("LINEBELOW", (0, -1), (-1, -1), 1.5, colors.HexColor("#000000")),
                            # Left border (thick)
                            ("LINELEFT", (0, 0), (0, -1), 1.5, colors.HexColor("#000000")),
                            # Right border (thick)
                            ("LINERIGHT", (-1, 0), (-1, -1), 1.5, colors.HexColor("#000000")),
                        ]
                    )
                )

                # Draw table
                table.wrapOn(c, available_width, available_height)
                table.drawOn(c, margin, page_height - 1 * inch - table.height)

            except Exception as table_error:
                logger.warning(f"Failed to create table, falling back to simple text: {table_error}")

                # Simple fallback: draw text directly
                y = page_height - 1 * inch
                c.setFont("Helvetica", cell_font_size)

                for row in data:
                    x = margin
                    col_width = available_width / len(row) if row else available_width

                    for cell_text in row:
                        # Truncate text if too long
                        truncated = cell_text[:20] if len(cell_text) > 20 else cell_text
                        c.drawString(x, y, truncated)
                        x += col_width

                    y -= 15
                    if y < margin:
                        c.showPage()
                        y = page_height - margin

        c.save()
        pdf_buffer.seek(0)

        logger.info("Successfully converted XLSX to PDF using openpyxl + reportlab")
        return pdf_buffer

    except ImportError as e:
        logger.error(f"Required libraries not installed for XLSX conversion: {e}")
        raise
    except Exception as e:
        logger.error(f"Failed to convert XLSX to PDF: {e}")
        raise
