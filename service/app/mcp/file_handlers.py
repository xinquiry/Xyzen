"""
File handlers for different file types using a Strategy pattern.
Supports reading and writing/creating files.
"""

import abc
import csv
import io
import logging
from typing import List, Literal, Union

from fastmcp import Image

logger = logging.getLogger(__name__)

ReadMode = Literal["text", "image"]


class BaseFileHandler(abc.ABC):
    """Abstract base handler for file operations."""

    @abc.abstractmethod
    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        """
        Reads content from file bytes.

        Args:
            file_bytes: The raw file content.
            mode: The reading mode ("text" or "image").

        Returns:
            Extracted text string or list of images.
        """
        pass

    @abc.abstractmethod
    def create_content(self, text_content: str) -> bytes:
        """
        Creates file bytes from text content.

        Args:
            text_content: The text to write to the file.

        Returns:
            The generated file bytes.
        """
        pass


class TextFileHandler(BaseFileHandler):
    """Handler for plain text and CSV files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        if mode == "image":
            raise ValueError("Text files cannot be read as images.")
        # Try UTF-8 first, then fall back to latin-1 or others if needed
        return file_bytes.decode("utf-8", errors="replace")

    def create_content(self, text_content: str) -> bytes:
        return text_content.encode("utf-8")


class PdfFileHandler(BaseFileHandler):
    """Handler for PDF files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) is required for PDF handling. Please install 'pymupdf'.")

        doc = fitz.open(stream=file_bytes, filetype="pdf")

        if mode == "image":
            images = []
            for page in doc:
                # Render page to image (pixmap) at 2x resolution for better quality
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                images.append(Image(data=img_bytes, format="png"))
            return images

        else:  # Text mode
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            return "\n\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) is required for PDF handling. Please install 'pymupdf'.")

        doc = fitz.open()
        page = doc.new_page()
        # Simple text insertion
        # For more complex layout, we'd need a more sophisticated approach
        # inserting at (50, 50) with default font
        page.insert_text((50, 50), text_content)

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()


class DocxFileHandler(BaseFileHandler):
    """Handler for Word (.docx) files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        if mode == "image":
            raise ValueError("DOCX files cannot be read as images (yet).")

        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX handling. Please install 'python-docx'.")

        buffer = io.BytesIO(file_bytes)
        doc = Document(buffer)

        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)

        return "\n".join(full_text)

    def create_content(self, text_content: str) -> bytes:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX handling. Please install 'python-docx'.")

        doc = Document()
        # Split by newlines to create paragraphs
        for line in text_content.split("\n"):
            doc.add_paragraph(line)

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()


class ExcelFileHandler(BaseFileHandler):
    """Handler for Excel (.xlsx) files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        if mode == "image":
            raise ValueError("Excel files cannot be read as images.")

        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for XLSX handling. Please install 'openpyxl'.")

        buffer = io.BytesIO(file_bytes)
        wb = openpyxl.load_workbook(buffer, read_only=True, data_only=True)

        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                # Filter out None values and join
                row_text = "\t".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for XLSX handling. Please install 'openpyxl'.")

        wb = openpyxl.Workbook()
        ws = wb.active
        if ws:
            # Naive CSV-like parsing: split by newlines for rows, commas/tabs for columns?
            # For simplicity, let's just put lines in first column associated with text
            # Or if the user provides CSV-like text, we try to parse it.

            # Let's try to interpret as CSV if possible, otherwise just lines
            reader = csv.reader(io.StringIO(text_content))
            for row in reader:
                if row:
                    ws.append(row)

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()


class PptxFileHandler(BaseFileHandler):
    """Handler for PowerPoint (.pptx) files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, List[Image]]:
        if mode == "image":
            raise ValueError("PPTX files cannot be read as images (yet).")

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        buffer = io.BytesIO(file_bytes)
        prs = Presentation(buffer)

        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Pyright doesn't know about shape.text dynamics well
                    text_parts.append(str(getattr(shape, "text")))

        return "\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        prs = Presentation()

        # Simple logic: create one slide with a text box containing all text
        # Or split by "--- Slide ---" if we want to get fancy?
        # For now, let's just make one slide with title and content layout

        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content

        # Split into title and body based on first newline
        lines = text_content.strip().split("\n", 1)

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        if lines:
            if title_shape:
                title_shape.text = lines[0]
            if len(lines) > 1 and body_shape:
                body_shape.text = lines[1]  # type: ignore

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()


class FileHandlerFactory:
    """Factory to get the appropriate file handler."""

    @staticmethod
    def get_handler(filename: str) -> BaseFileHandler:
        filename_lower = filename.lower()

        if filename_lower.endswith(".pdf"):
            return PdfFileHandler()
        elif filename_lower.endswith(".docx") or filename_lower.endswith(".doc"):
            return DocxFileHandler()
        elif filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls"):
            return ExcelFileHandler()
        elif filename_lower.endswith(".pptx") or filename_lower.endswith(".ppt"):
            return PptxFileHandler()
        else:
            # Default to text handler for .txt, .csv, .md, .py, etc.
            return TextFileHandler()
