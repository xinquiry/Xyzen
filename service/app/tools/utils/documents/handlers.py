"""
File handlers for different file types using a Strategy pattern.
Supports reading and writing/creating files.

Handlers support both plain text input (backward compatible) and
structured JSON input via DocumentSpec/SpreadsheetSpec/PresentationSpec
for production-ready document generation.
"""

from __future__ import annotations

import abc
import csv
import io
import json
import logging
from typing import TYPE_CHECKING, Any, Literal, Union

from fastmcp.utilities.types import Image
from pydantic import BaseModel, ValidationError

if TYPE_CHECKING:
    from app.tools.utils.documents.spec import DocumentSpec, PresentationSpec, SpreadsheetSpec

logger = logging.getLogger(__name__)

ReadMode = Literal["text", "image"]


class BaseFileHandler(abc.ABC):
    """Abstract base handler for file operations."""

    @abc.abstractmethod
    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        """
        Reads content from file bytes.

        Args:
            file_bytes: The raw file content.
            mode: The reading mode ("text" or "image").

        Returns:
            Extracted text string or list of images.
        """
        pass

    @abc.abstractmethod
    def create_content(self, text_content: str) -> bytes:
        """
        Creates file bytes from text content.

        Args:
            text_content: The text to write to the file.
                Can be plain text or JSON string of a document spec.

        Returns:
            The generated file bytes.
        """
        pass

    def _try_parse_spec(self, text_content: str, spec_class: type[BaseModel]) -> BaseModel | None:
        """
        Try to parse text content as a document specification JSON.

        Args:
            text_content: The input string (may be plain text or JSON).
            spec_class: The Pydantic model class to validate against.

        Returns:
            Parsed spec model if valid JSON matching the schema, None otherwise.
        """
        try:
            data = json.loads(text_content)
            if isinstance(data, dict):
                return spec_class.model_validate(data)
        except (json.JSONDecodeError, ValidationError):
            pass
        return None


class TextFileHandler(BaseFileHandler):
    """Handler for plain text files (.txt, .md, .csv, etc.)."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("Text files cannot be read as images.")
        return file_bytes.decode("utf-8", errors="replace")

    def create_content(self, text_content: str) -> bytes:
        return text_content.encode("utf-8")


class HtmlFileHandler(BaseFileHandler):
    """Handler for HTML files with BeautifulSoup text extraction."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("HTML files cannot be read as images.")

        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 is required for HTML handling. Please install 'beautifulsoup4'.")

        html_content = file_bytes.decode("utf-8", errors="replace")
        soup = BeautifulSoup(html_content, "lxml")

        # Remove script and style elements
        for element in soup(["script", "style", "meta", "link", "noscript"]):
            element.decompose()

        # Extract text with some structure preservation
        text_parts: list[str] = []

        # Process headings
        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
            level = int(tag.name[1])
            prefix = "#" * level
            text_parts.append(f"{prefix} {tag.get_text(strip=True)}")

        # Get all text if no structured extraction worked
        if not text_parts:
            text = soup.get_text(separator="\n", strip=True)
            # Clean up excessive whitespace
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            return "\n".join(lines)

        return "\n\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        """Create a simple HTML document from text."""
        # Escape HTML entities
        escaped = text_content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

        # Convert newlines to paragraphs
        paragraphs = escaped.split("\n\n")
        body_content = "\n".join(f"<p>{p.replace(chr(10), '<br>')}</p>" for p in paragraphs if p.strip())

        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Document</title>
</head>
<body>
{body_content}
</body>
</html>"""
        return html.encode("utf-8")


class JsonFileHandler(BaseFileHandler):
    """Handler for JSON files with pretty-printing."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("JSON files cannot be read as images.")

        text = file_bytes.decode("utf-8", errors="replace")
        try:
            data = json.loads(text)
            # Pretty-print for readability
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            # Return as-is if not valid JSON
            return text

    def create_content(self, text_content: str) -> bytes:
        """Create JSON file, validating and formatting if possible."""
        try:
            # Try to parse and re-format
            data = json.loads(text_content)
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            return formatted.encode("utf-8")
        except json.JSONDecodeError:
            # If not valid JSON, wrap as string value
            wrapped = json.dumps({"content": text_content}, indent=2, ensure_ascii=False)
            return wrapped.encode("utf-8")


class YamlFileHandler(BaseFileHandler):
    """Handler for YAML files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("YAML files cannot be read as images.")

        try:
            import yaml
        except ImportError:
            raise ImportError("pyyaml is required for YAML handling. Please install 'pyyaml'.")

        text = file_bytes.decode("utf-8", errors="replace")
        try:
            data = yaml.safe_load(text)
            # Pretty-print as YAML
            return yaml.dump(data, default_flow_style=False, allow_unicode=True)
        except yaml.YAMLError:
            return text

    def create_content(self, text_content: str) -> bytes:
        """Create YAML file from text or JSON."""
        try:
            import yaml
        except ImportError:
            raise ImportError("pyyaml is required for YAML handling. Please install 'pyyaml'.")

        try:
            # Try to parse as JSON first
            data = json.loads(text_content)
            return yaml.dump(data, default_flow_style=False, allow_unicode=True).encode("utf-8")
        except json.JSONDecodeError:
            try:
                # Try as YAML
                data = yaml.safe_load(text_content)
                return yaml.dump(data, default_flow_style=False, allow_unicode=True).encode("utf-8")
            except yaml.YAMLError:
                # Return as plain text
                return text_content.encode("utf-8")


class XmlFileHandler(BaseFileHandler):
    """Handler for XML files."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("XML files cannot be read as images.")

        import xml.etree.ElementTree as ET

        text = file_bytes.decode("utf-8", errors="replace")
        try:
            root = ET.fromstring(text)
            # Extract text content recursively
            return self._extract_xml_text(root)
        except ET.ParseError:
            return text

    def _extract_xml_text(self, element: Any, depth: int = 0) -> str:
        """Recursively extract text from XML with structure hints."""
        parts: list[str] = []
        indent = "  " * depth

        tag_name = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        # Add element text
        if element.text and element.text.strip():
            parts.append(f"{indent}[{tag_name}] {element.text.strip()}")

        # Process children
        for child in element:
            parts.append(self._extract_xml_text(child, depth + 1))

        # Add tail text
        if element.tail and element.tail.strip():
            parts.append(f"{indent}{element.tail.strip()}")

        return "\n".join(parts)

    def create_content(self, text_content: str) -> bytes:
        """Create a simple XML document from text."""
        import xml.etree.ElementTree as ET

        # Escape XML entities
        escaped = text_content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

        root = ET.Element("document")
        content = ET.SubElement(root, "content")
        content.text = escaped

        return ET.tostring(root, encoding="unicode", xml_declaration=True).encode("utf-8")


class ImageFileHandler(BaseFileHandler):
    """Handler for image files with OCR support."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            # Return the image as-is
            format_type = self._detect_format(file_bytes)
            return [Image(data=file_bytes, format=format_type)]

        # Text mode - use OCR
        try:
            from PIL import Image as PILImage
        except ImportError:
            raise ImportError("pillow is required for image handling. Please install 'pillow'.")

        try:
            import pytesseract
        except ImportError:
            raise ImportError(
                "pytesseract is required for OCR. Please install 'pytesseract' and Tesseract OCR system package."
            )

        img = PILImage.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img)
        return text.strip() if text else "(No text detected in image)"

    def _detect_format(self, file_bytes: bytes) -> str:
        """Detect image format from magic bytes."""
        if file_bytes[:8] == b"\x89PNG\r\n\x1a\n":
            return "png"
        elif file_bytes[:2] == b"\xff\xd8":
            return "jpeg"
        elif file_bytes[:6] in (b"GIF87a", b"GIF89a"):
            return "gif"
        elif file_bytes[:4] == b"RIFF" and file_bytes[8:12] == b"WEBP":
            return "webp"
        elif file_bytes[:2] in (b"BM", b"BA"):
            return "bmp"
        else:
            return "png"  # Default

    def create_content(self, text_content: str) -> bytes:
        """Images cannot be created from text."""
        raise ValueError("Cannot create image files from text content. Use an image editing tool instead.")


class PdfFileHandler(BaseFileHandler):
    """Handler for PDF files with table extraction and production-ready generation."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) is required for PDF handling. Please install 'pymupdf'.")

        doc = fitz.open(stream=file_bytes, filetype="pdf")

        if mode == "image":
            images = []
            for page in doc:
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                images.append(Image(data=img_bytes, format="png"))
            return images

        # Text mode with table extraction
        text_parts: list[str] = []
        for page in doc:
            # Try table extraction first
            try:
                tables = page.find_tables()  # type: ignore[attr-defined]
                if tables and tables.tables:
                    for table in tables.tables:
                        text_parts.append(self._format_table(table))
            except Exception:
                pass  # Table extraction not available or failed

            # Get text with layout preservation
            text: str = page.get_text("text", sort=True)  # type: ignore[attr-defined]
            if text.strip():
                text_parts.append(text)

        return "\n\n".join(text_parts)

    def _format_table(self, table: Any) -> str:
        """Format a PyMuPDF table as text."""
        rows = table.extract()
        if not rows:
            return ""

        # Calculate column widths
        col_widths: list[int] = []
        for row in rows:
            for i, cell in enumerate(row):
                cell_str = str(cell) if cell else ""
                if i >= len(col_widths):
                    col_widths.append(len(cell_str))
                else:
                    col_widths[i] = max(col_widths[i], len(cell_str))

        # Format rows
        lines: list[str] = []
        for i, row in enumerate(rows):
            cells = [str(cell if cell else "").ljust(col_widths[j]) for j, cell in enumerate(row)]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:  # Add separator after header
                lines.append("|" + "|".join("-" * (w + 2) for w in col_widths) + "|")

        return "\n".join(lines)

    def create_content(self, text_content: str) -> bytes:
        """Create PDF from text or DocumentSpec JSON."""
        from app.tools.utils.documents.spec import DocumentSpec

        spec = self._try_parse_spec(text_content, DocumentSpec)
        if spec:
            return self._create_pdf_from_spec(spec)  # type: ignore[arg-type]
        return self._create_pdf_from_text(text_content)

    def _create_pdf_from_text(self, text_content: str) -> bytes:
        """Create simple PDF from plain text using reportlab."""
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story: list[Any] = []

        for line in text_content.split("\n"):
            if line.strip():
                story.append(Paragraph(line, styles["Normal"]))
                story.append(Spacer(1, 12))

        if not story:
            story.append(Paragraph("(Empty document)", styles["Normal"]))

        doc.build(story)
        return buffer.getvalue()

    def _create_pdf_from_spec(self, spec: DocumentSpec) -> bytes:
        """Create production PDF from DocumentSpec."""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4, legal, letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import (
            PageBreak,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )

        buffer = io.BytesIO()
        page_sizes = {"letter": letter, "A4": A4, "legal": legal}
        page_size = page_sizes.get(spec.page_size, letter)

        doc = SimpleDocTemplate(
            buffer,
            pagesize=page_size,
            title=spec.title or "",
            author=spec.author or "",
            subject=spec.subject or "",
        )

        styles = getSampleStyleSheet()
        story: list[Any] = []

        for block in spec.content:
            if block.type == "heading":
                level = min(block.level, 6)  # type: ignore[union-attr]
                style_name = f"Heading{level}" if level <= 3 else "Heading3"
                story.append(Paragraph(block.content, styles[style_name]))  # type: ignore[union-attr]
            elif block.type == "text":
                style = styles["Normal"]
                if hasattr(block, "style") and block.style == "code":  # type: ignore[union-attr]
                    style = styles["Code"]
                story.append(Paragraph(block.content, style))  # type: ignore[union-attr]
            elif block.type == "list":
                list_items = [
                    Paragraph(f"{'• ' if not block.ordered else f'{i + 1}. '}{item}", styles["Normal"])
                    for i, item in enumerate(block.items)
                ]
                for item in list_items:
                    story.append(item)
            elif block.type == "table":
                table_data = [block.headers] + block.rows  # type: ignore[union-attr]
                t = Table(table_data)
                t.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                            ("FONTSIZE", (0, 0), (-1, 0), 11),
                            ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                            ("GRID", (0, 0), (-1, -1), 1, colors.black),
                        ]
                    )
                )
                story.append(t)
            elif block.type == "page_break":
                story.append(PageBreak())

            story.append(Spacer(1, 12))

        if not story:
            story.append(Paragraph("(Empty document)", styles["Normal"]))

        doc.build(story)
        return buffer.getvalue()


class DocxFileHandler(BaseFileHandler):
    """Handler for Word (.docx) files with table support."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("DOCX files cannot be read as images.")

        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX handling. Please install 'python-docx'.")

        buffer = io.BytesIO(file_bytes)
        doc = Document(buffer)

        full_text: list[str] = []

        # Process document body elements in order
        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                # Paragraph
                para_text = "".join(node.text or "" for node in element.iter() if node.text)
                if para_text.strip():
                    full_text.append(para_text)
            elif tag == "tbl":
                # Table - extract and format
                table_text = self._extract_table(element, doc)
                if table_text:
                    full_text.append(table_text)

        return "\n\n".join(full_text)

    def _extract_table(self, tbl_element: Any, doc: Any) -> str:
        """Extract table from DOCX element."""
        try:
            from docx.table import Table
        except ImportError:
            return ""

        table = Table(tbl_element, doc)
        rows: list[list[str]] = []

        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        # Calculate column widths
        col_widths = [max(len(row[i]) if i < len(row) else 0 for row in rows) for i in range(len(rows[0]))]

        # Format as text table
        lines: list[str] = []
        for i, row in enumerate(rows):
            cells = [cell.ljust(col_widths[j]) for j, cell in enumerate(row)]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + "|".join("-" * (w + 2) for w in col_widths) + "|")

        return "\n".join(lines)

    def create_content(self, text_content: str) -> bytes:
        """Create DOCX from text or DocumentSpec JSON."""
        from app.tools.utils.documents.spec import DocumentSpec

        spec = self._try_parse_spec(text_content, DocumentSpec)
        if spec:
            return self._create_docx_from_spec(spec)  # type: ignore[arg-type]
        return self._create_docx_from_text(text_content)

    def _create_docx_from_text(self, text_content: str) -> bytes:
        """Create simple DOCX from plain text."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX handling. Please install 'python-docx'.")

        doc = Document()
        for line in text_content.split("\n"):
            doc.add_paragraph(line)

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _create_docx_from_spec(self, spec: DocumentSpec) -> bytes:
        """Create production DOCX from DocumentSpec."""
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            raise ImportError("python-docx is required for DOCX handling. Please install 'python-docx'.")

        doc = Document()

        # Set document properties
        if spec.title:
            doc.core_properties.title = spec.title
        if spec.author:
            doc.core_properties.author = spec.author
        if spec.subject:
            doc.core_properties.subject = spec.subject

        for block in spec.content:
            if block.type == "heading":
                level = min(block.level, 9)  # type: ignore[union-attr]
                doc.add_heading(block.content, level=level)  # type: ignore[union-attr]
            elif block.type == "text":
                p = doc.add_paragraph(block.content)  # type: ignore[union-attr]
                if hasattr(block, "style"):
                    if block.style == "bold" and p.runs:  # type: ignore[union-attr]
                        p.runs[0].bold = True
                    elif block.style == "italic" and p.runs:  # type: ignore[union-attr]
                        p.runs[0].italic = True
                    elif block.style == "code" and p.runs:  # type: ignore[union-attr]
                        p.runs[0].font.name = "Courier New"
                        p.runs[0].font.size = Pt(10)
            elif block.type == "list":
                style = "List Number" if block.ordered else "List Bullet"  # type: ignore[union-attr]
                for item in block.items:  # type: ignore[union-attr]
                    doc.add_paragraph(item, style=style)
            elif block.type == "table":
                table = doc.add_table(rows=1, cols=len(block.headers))  # type: ignore[union-attr]
                table.style = "Table Grid"
                # Header row
                hdr_cells = table.rows[0].cells
                for i, header in enumerate(block.headers):  # type: ignore[union-attr]
                    hdr_cells[i].text = header
                # Data rows
                for row_data in block.rows:  # type: ignore[union-attr]
                    row_cells = table.add_row().cells
                    for i, cell_val in enumerate(row_data):
                        row_cells[i].text = str(cell_val)
            elif block.type == "page_break":
                doc.add_page_break()

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()


class ExcelFileHandler(BaseFileHandler):
    """Handler for Excel (.xlsx) files with formatting support."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("Excel files cannot be read as images.")

        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for XLSX handling. Please install 'openpyxl'.")

        buffer = io.BytesIO(file_bytes)
        wb = openpyxl.load_workbook(buffer, read_only=True, data_only=True)

        text_parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        """Create XLSX from text or SpreadsheetSpec JSON."""
        from app.tools.utils.documents.spec import SpreadsheetSpec

        spec = self._try_parse_spec(text_content, SpreadsheetSpec)
        if spec:
            return self._create_xlsx_from_spec(spec)  # type: ignore[arg-type]
        return self._create_xlsx_from_csv(text_content)

    def _create_xlsx_from_csv(self, text_content: str) -> bytes:
        """Create simple XLSX from CSV-like text."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for XLSX handling. Please install 'openpyxl'.")

        wb = openpyxl.Workbook()
        ws = wb.active
        if ws:
            reader = csv.reader(io.StringIO(text_content))
            for row in reader:
                if row:
                    ws.append(row)

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    def _create_xlsx_from_spec(self, spec: SpreadsheetSpec) -> bytes:
        """Create production XLSX from SpreadsheetSpec."""
        try:
            import openpyxl
            from openpyxl.styles import Border, Font, PatternFill, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            raise ImportError("openpyxl is required for XLSX handling. Please install 'openpyxl'.")

        wb = openpyxl.Workbook()
        # Remove default sheet if we have sheets to add
        if spec.sheets and wb.active:
            wb.remove(wb.active)

        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        for sheet_spec in spec.sheets:
            ws = wb.create_sheet(title=sheet_spec.name)

            start_row = 1
            num_cols = (
                len(sheet_spec.headers) if sheet_spec.headers else (len(sheet_spec.data[0]) if sheet_spec.data else 1)
            )

            # Add headers if provided
            if sheet_spec.headers:
                for col, header in enumerate(sheet_spec.headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.border = thin_border
                start_row = 2
                if sheet_spec.freeze_header:
                    ws.freeze_panes = "A2"

            # Add data
            for row_idx, row_data in enumerate(sheet_spec.data, start_row):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = thin_border

            # Auto-adjust column widths
            for col_idx in range(1, num_cols + 1):
                col_letter = get_column_letter(col_idx)
                max_length = 0
                for row_idx in range(1, len(sheet_spec.data) + start_row):
                    cell_value = ws.cell(row=row_idx, column=col_idx).value
                    if cell_value:
                        max_length = max(max_length, len(str(cell_value)))
                ws.column_dimensions[col_letter].width = min(max_length + 2, 50)

        # Ensure at least one sheet exists
        if not wb.sheetnames:
            wb.create_sheet("Sheet1")

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()


class PptxFileHandler(BaseFileHandler):
    """Handler for PowerPoint (.pptx) files with multi-slide support."""

    def read_content(self, file_bytes: bytes, mode: ReadMode = "text") -> Union[str, list[Image]]:
        if mode == "image":
            raise ValueError("PPTX files cannot be read as images.")

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        buffer = io.BytesIO(file_bytes)
        prs = Presentation(buffer)

        text_parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:  # type: ignore[union-attr]
                    text_parts.append(str(shape.text))  # type: ignore[union-attr]

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    text_parts.append(f"[Notes: {notes.strip()}]")

        return "\n".join(text_parts)

    def create_content(self, text_content: str) -> bytes:
        """Create PPTX from text or PresentationSpec JSON."""
        from app.tools.utils.documents.spec import PresentationSpec

        spec = self._try_parse_spec(text_content, PresentationSpec)
        if spec:
            return self._create_pptx_from_spec(spec)  # type: ignore[arg-type]
        return self._create_pptx_from_text(text_content)

    def _create_pptx_from_text(self, text_content: str) -> bytes:
        """Create simple PPTX from plain text."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        prs = Presentation()
        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content

        lines = text_content.strip().split("\n", 1)

        slide = prs.slides.add_slide(bullet_slide_layout)
        if slide.shapes.title and lines:
            slide.shapes.title.text = lines[0]
        if len(lines) > 1 and len(slide.placeholders) > 1:
            slide.placeholders[1].text = lines[1]  # type: ignore[union-attr]

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _create_pptx_from_spec(self, spec: PresentationSpec) -> bytes:
        """Create production PPTX from PresentationSpec."""
        # Route based on mode
        if spec.mode == "image_slides":
            return self._create_pptx_image_slides(spec)
        else:
            return self._create_pptx_structured(spec)

    def _create_pptx_structured(self, spec: PresentationSpec) -> bytes:
        """Create PPTX with structured DSL slides (traditional mode)."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        from app.tools.utils.documents.image_fetcher import ImageFetcher

        prs = Presentation()
        image_fetcher = ImageFetcher()

        # Layout mapping
        LAYOUTS = {
            "title": 0,
            "title_content": 1,
            "section": 2,
            "two_column": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6,
        }

        for slide_spec in spec.slides:
            layout_idx = LAYOUTS.get(slide_spec.layout, 1)
            # Ensure layout index is valid
            layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if slide_spec.title and slide.shapes.title:
                slide.shapes.title.text = slide_spec.title

            # Set subtitle (for title slides)
            if slide_spec.subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_spec.subtitle  # type: ignore[union-attr]

            # Render content blocks
            if slide_spec.content:
                self._render_content_blocks(slide, slide_spec.content, image_fetcher)

            # Add speaker notes
            if slide_spec.notes:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_slide.notes_text_frame.text = slide_spec.notes

        # Ensure at least one slide
        if not prs.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _create_pptx_image_slides(self, spec: PresentationSpec) -> bytes:
        """Create PPTX with full-bleed images as slides."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx is required for PPTX handling. Please install 'python-pptx'.")

        from app.tools.utils.documents.image_fetcher import ImageFetcher

        prs = Presentation()
        # Set slide dimensions (16:9 widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        image_fetcher = ImageFetcher()
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for slide_spec in spec.image_slides:
            slide = prs.slides.add_slide(blank_layout)

            # Use storage_url if available (resolved by async layer), otherwise fall back to image_id
            if slide_spec.storage_url:
                result = image_fetcher.fetch(url=slide_spec.storage_url)
            else:
                result = image_fetcher.fetch(image_id=slide_spec.image_id)

            if result.success and result.data:
                # Add full-bleed image (0,0 to full slide dimensions)
                image_stream = io.BytesIO(result.data)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(0),
                    Inches(0),
                    prs.slide_width,
                    prs.slide_height,
                )
            else:
                # Add error text for failed images
                text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
                tf = text_box.text_frame
                tf.paragraphs[0].text = f"[Slide image failed: {result.error}]"
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.italic = True

            # Add speaker notes
            if slide_spec.notes:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_slide.notes_text_frame.text = slide_spec.notes

        # Ensure at least one slide exists
        if not prs.slides:
            slide = prs.slides.add_slide(blank_layout)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _render_content_blocks(
        self,
        slide: Any,
        content_blocks: list[Any],
        image_fetcher: Any,
    ) -> None:
        """Render all content blocks on a slide with vertical stacking."""

        # Content area dimensions (below title)
        CONTENT_LEFT = 0.5  # inches
        CONTENT_TOP = 1.8  # inches
        CONTENT_WIDTH = 9.0  # inches
        CONTENT_BOTTOM = 7.0  # inches

        current_y = CONTENT_TOP

        for block in content_blocks:
            if current_y >= CONTENT_BOTTOM:
                logger.warning("Slide content area full, skipping remaining blocks")
                break

            remaining_height = CONTENT_BOTTOM - current_y

            if block.type == "text":
                height = self._render_text_block(slide, block, CONTENT_LEFT, current_y, CONTENT_WIDTH)
            elif block.type == "list":
                height = self._render_list_block(slide, block, CONTENT_LEFT, current_y, CONTENT_WIDTH)
            elif block.type == "image":
                height = self._render_image_block(
                    slide, block, CONTENT_LEFT, current_y, CONTENT_WIDTH, remaining_height, image_fetcher
                )
            elif block.type == "table":
                height = self._render_table_block(
                    slide, block, CONTENT_LEFT, current_y, CONTENT_WIDTH, remaining_height
                )
            elif block.type == "heading":
                height = self._render_heading_block(slide, block, CONTENT_LEFT, current_y, CONTENT_WIDTH)
            else:
                # Unknown block type, skip
                height = 0.0

            current_y += height

    def _render_text_block(
        self,
        slide: Any,
        block: Any,
        left: float,
        top: float,
        max_width: float,
    ) -> float:
        """Render a text block. Returns height in inches."""
        from pptx.util import Inches, Pt

        # Estimate height based on text length
        chars_per_line = int(max_width * 12)  # ~12 chars per inch at 12pt
        num_lines = max(1, len(block.content) // chars_per_line + 1)
        box_height = num_lines * 0.25  # ~0.25 inches per line

        text_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(max_width),
            Inches(box_height),
        )

        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = block.content
        p.font.size = Pt(12)

        # Apply style
        if hasattr(block, "style"):
            if block.style == "bold":
                p.font.bold = True
            elif block.style == "italic":
                p.font.italic = True
            elif block.style == "code":
                p.font.name = "Courier New"
                p.font.size = Pt(10)

        return box_height + 0.1  # Add margin

    def _render_list_block(
        self,
        slide: Any,
        block: Any,
        left: float,
        top: float,
        max_width: float,
    ) -> float:
        """Render a list block. Returns height in inches."""
        from pptx.util import Inches, Pt

        num_items = len(block.items)
        item_height = 0.3  # inches per item
        box_height = num_items * item_height

        text_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(max_width),
            Inches(box_height),
        )

        tf = text_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(block.items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = f"{i + 1}. " if block.ordered else "• "
            p.text = prefix + item
            p.font.size = Pt(12)
            p.level = 0

        return box_height + 0.1

    def _render_heading_block(
        self,
        slide: Any,
        block: Any,
        left: float,
        top: float,
        max_width: float,
    ) -> float:
        """Render a heading block. Returns height in inches."""
        from pptx.util import Inches, Pt

        # Font sizes by heading level
        HEADING_SIZES = {
            1: 24,
            2: 20,
            3: 18,
            4: 16,
            5: 14,
            6: 12,
        }

        level = getattr(block, "level", 1)
        font_size = HEADING_SIZES.get(level, 14)
        box_height = font_size / 72.0 * 1.5  # Convert to inches with padding

        text_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(max_width),
            Inches(box_height),
        )

        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = block.content
        p.font.size = Pt(font_size)
        p.font.bold = True

        return box_height + 0.1

    def _render_table_block(
        self,
        slide: Any,
        block: Any,
        left: float,
        top: float,
        max_width: float,
        max_height: float,
    ) -> float:
        """Render a table block. Returns height in inches."""
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        num_cols = len(block.headers)
        if num_cols == 0:
            return 0.0

        num_rows = 1 + len(block.rows)  # Header + data rows

        # Calculate row height
        row_height = 0.4  # inches
        table_height = num_rows * row_height

        # Cap table height to available space
        max_table_height = min(4.0, max_height - 0.2)
        if table_height > max_table_height:
            table_height = max_table_height
            row_height = table_height / num_rows

        # Create table
        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            Inches(left),
            Inches(top),
            Inches(max_width),
            Inches(table_height),
        )
        table = table_shape.table

        # Style header row
        for i, header in enumerate(block.headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

            # Set text properties
            if cell.text_frame.paragraphs:
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(11)
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para.alignment = PP_ALIGN.CENTER

        # Fill data rows
        for row_idx, row_data in enumerate(block.rows):
            for col_idx, cell_val in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_val)
                    if cell.text_frame.paragraphs:
                        para = cell.text_frame.paragraphs[0]
                        para.font.size = Pt(10)

        return table_height + 0.2

    def _render_image_block(
        self,
        slide: Any,
        block: Any,
        left: float,
        top: float,
        max_width: float,
        max_height: float,
        image_fetcher: Any,
    ) -> float:
        """Render an image block. Returns height in inches."""
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        # Fetch image by url or image_id
        url = getattr(block, "url", None)
        image_id = getattr(block, "image_id", None)
        result = image_fetcher.fetch(url=url, image_id=image_id)

        if not result.success:
            # Add error placeholder
            text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(max_width), Inches(0.5))
            tf = text_box.text_frame
            tf.paragraphs[0].text = f"[Image failed to load: {result.error}]"
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].font.size = Pt(10)
            return 0.6

        # Calculate image dimensions
        if block.width:
            # Use specified width (in points, convert to inches)
            img_width = block.width / 72.0
        elif result.width and result.height:
            # Scale to fit max_width while maintaining aspect ratio
            img_width = min(max_width * 0.8, result.width / 96.0)  # 96 DPI assumption, 80% max width
        else:
            img_width = min(max_width * 0.6, 4.0)  # Default 4 inches or 60% width

        # Calculate height maintaining aspect ratio
        if result.width and result.height:
            aspect = result.height / result.width
            img_height = img_width * aspect
        else:
            img_height = img_width * 0.75  # Default 4:3 aspect

        # Cap height to available space (leave room for caption)
        caption_space = 0.5 if block.caption else 0.1
        available_height = max_height - caption_space
        if img_height > available_height:
            scale = available_height / img_height
            img_height = available_height
            img_width = img_width * scale

        # Center image horizontally
        img_left = left + (max_width - img_width) / 2

        # Add image to slide
        image_stream = io.BytesIO(result.data)
        slide.shapes.add_picture(
            image_stream,
            Inches(img_left),
            Inches(top),
            Inches(img_width),
            Inches(img_height),
        )

        total_height = img_height + 0.1

        # Add caption if present
        if block.caption:
            caption_top = top + img_height + 0.1
            caption_box = slide.shapes.add_textbox(Inches(left), Inches(caption_top), Inches(max_width), Inches(0.3))
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = block.caption
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.italic = True
            total_height += 0.4

        return total_height


class FileHandlerFactory:
    """Factory to get the appropriate file handler based on filename."""

    @staticmethod
    def get_handler(filename: str) -> BaseFileHandler:
        filename_lower = filename.lower()

        # HTML files
        if filename_lower.endswith((".html", ".htm")):
            return HtmlFileHandler()

        # Data interchange formats
        elif filename_lower.endswith(".json"):
            return JsonFileHandler()
        elif filename_lower.endswith((".yaml", ".yml")):
            return YamlFileHandler()
        elif filename_lower.endswith(".xml"):
            return XmlFileHandler()

        # Image files
        elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp")):
            return ImageFileHandler()

        # Office documents
        elif filename_lower.endswith(".pdf"):
            return PdfFileHandler()
        elif filename_lower.endswith((".docx", ".doc")):
            return DocxFileHandler()
        elif filename_lower.endswith((".xlsx", ".xls")):
            return ExcelFileHandler()
        elif filename_lower.endswith((".pptx", ".ppt")):
            return PptxFileHandler()

        # Default to text handler
        else:
            return TextFileHandler()


__all__ = [
    "BaseFileHandler",
    "TextFileHandler",
    "HtmlFileHandler",
    "JsonFileHandler",
    "YamlFileHandler",
    "XmlFileHandler",
    "ImageFileHandler",
    "PdfFileHandler",
    "DocxFileHandler",
    "ExcelFileHandler",
    "PptxFileHandler",
    "FileHandlerFactory",
    "ReadMode",
]
